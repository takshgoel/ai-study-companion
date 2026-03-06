import io
from hashlib import sha1

import PyPDF2
from pptx import Presentation

SUPPORTED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}


def human_file_size(num_bytes: int) -> str:
    if num_bytes < 1024:
        return f"{num_bytes} B"
    if num_bytes < 1024 * 1024:
        return f"{num_bytes / 1024:.1f} KB"
    return f"{num_bytes / (1024 * 1024):.1f} MB"


def file_id(name: str, data: bytes) -> str:
    return sha1(name.encode("utf-8") + data).hexdigest()[:16]


def extract_text_from_pdf_bytes(data: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(data))
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)

    return "\n".join(text_parts).strip()


def extract_text_from_pptx_bytes(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    text_parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

    return "\n".join(text_parts).strip()


def parse_uploaded_file(uploaded_file) -> dict:
    data = uploaded_file.getvalue()
    ext = SUPPORTED_TYPES.get(uploaded_file.type)

    if ext is None:
        return {
            "ok": False,
            "error": f"Unsupported file type for {uploaded_file.name}.",
        }

    parsed_text = ""

    try:
        if ext == "pdf":
            parsed_text = extract_text_from_pdf_bytes(data)
        else:
            parsed_text = extract_text_from_pptx_bytes(data)
    except Exception as exc:
        return {
            "ok": False,
            "error": f"Could not parse {uploaded_file.name}: {exc}",
        }

    if not parsed_text.strip():
        return {
            "ok": False,
            "error": f"No text could be extracted from {uploaded_file.name}.",
        }

    return {
        "ok": True,
        "id": file_id(uploaded_file.name, data),
        "name": uploaded_file.name,
        "type": uploaded_file.type,
        "size": len(data),
        "size_human": human_file_size(len(data)),
        "content": parsed_text,
    }
